from pathlib import Path
from typing import Literal
import pandas as pd
from PIL import Image
import cv2
import numpy as np
import pytesseract

DocType = Literal["pdf", "docx", "pptx", "txt", "csv", "xlsx", "image", "unknown"]

_SUFFIX_MAP: dict[str, DocType] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".csv": "csv",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
}


def detect_type(path: Path) -> DocType:
    return _SUFFIX_MAP.get(path.suffix.lower(), "unknown")


def extract_text(path: Path, kind: DocType) -> str:
    if kind == "pdf":
        from pypdf import PdfReader

        text = []
        reader = PdfReader(str(path))
        for page in reader.pages:
            text.append(page.extract_text() or "")
        return "\n".join(text).strip()

    if kind == "docx":
        import docx2txt

        return (docx2txt.process(str(path)) or "").strip()

    if kind == "pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out).strip()

    if kind == "txt":
        return path.read_text(encoding="utf-8", errors="replace")

    if kind == "csv":
        df = pd.read_csv(path)
        return _df_to_text(df)

    if kind == "xlsx":
        df = pd.read_excel(path, engine="openpyxl")
        return _df_to_text(df)

    if kind == "image":
        return _ocr_image(path)

    return ""


def _df_to_text(df: pd.DataFrame) -> str:
    lines = []
    lines.append(" | ".join(map(str, df.columns.tolist())))
    for _, row in df.head(100).iterrows():
        lines.append(" | ".join(map(lambda x: str(x), row.tolist())))
    return "\n".join(lines)


def _ocr_image(path: Path) -> str:
    # Load with cv2 for preprocessing
    img = cv2.imdecode(np.fromfile(str(path), dtype=np.uint8), cv2.IMREAD_COLOR)
    if img is None:
        # Fallback via PIL if cv2 fails
        pil = Image.open(path)
        return pytesseract.image_to_string(pil)

    # Preprocess: grayscale -> threshold -> slight dilation to connect text
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    # Adaptive threshold handles uneven lighting
    thr = cv2.adaptiveThreshold(
        gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 11
    )
    # Optional morphological ops (light)
    kernel = np.ones((1, 1), np.uint8)
    proc = cv2.morphologyEx(thr, cv2.MORPH_OPEN, kernel)

    # Convert back to PIL for pytesseract
    pil = Image.fromarray(proc)
    text = pytesseract.image_to_string(pil, config="--psm 6")  # Assume block of text
    return text.strip()
